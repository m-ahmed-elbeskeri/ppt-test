"""Deterministic ground truth via python-pptx.

This is the suite's measuring stick: a reproducible extraction of the text,
tables, speaker notes, chart data and image references that *exist* in a deck.
Metric denominators (text/notes/table/image recall) are computed against it.

HONEST LIMITATION (surfaced in the report): python-pptx cannot see text inside
SmartArt that lacks a text frame, text rasterised into images (EMF/WDP/PNG), or
some exotic embedded objects. So recall is "recall vs python-pptx-extractable
content" — a consistent *relative* yardstick across converters, not absolute
truth. Rendered slides + human ratings are the corrective for what GT misses.
"""
from __future__ import annotations

from pathlib import Path

from .config import deck_slug
from .schema import GroundTruth, SlideGT
from .tokens import count_tokens


def _walk(shapes):
    """Yield every shape, recursing into groups."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for sh in shapes:
        yield sh
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _walk(sh.shapes)
        except Exception:
            # Some shapes raise on shape_type; treat as leaf.
            continue


def _chart_text(shape) -> list[str]:
    """Best-effort: pull category labels + series names/values from a chart so the
    deck's *data* counts as ground-truth content."""
    out: list[str] = []
    try:
        chart = shape.chart
    except Exception:
        return out
    try:
        cats = [str(c) for c in chart.plots[0].categories if c is not None]
        out.extend(cats)
    except Exception:
        pass
    try:
        for series in chart.series:
            try:
                if series.name:
                    out.append(str(series.name))
            except Exception:
                pass
            try:
                out.extend(str(v) for v in series.values if v is not None)
            except Exception:
                pass
    except Exception:
        pass
    return out


def _slide_images(slide) -> list[str]:
    """Image media filenames referenced by a slide (via its relationships)."""
    names: list[str] = []
    try:
        for rel in slide.part.rels.values():
            try:
                if rel.is_external:
                    continue
                if "image" in rel.reltype:
                    partname = str(rel.target_part.partname)
                    names.append(partname.split("/")[-1])
            except Exception:
                continue
    except Exception:
        pass
    return names


def extract_slide(slide, number: int) -> SlideGT:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    title = None
    try:
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            t = slide.shapes.title.text.strip()
            title = t or None
    except Exception:
        title = None

    text_runs: list[str] = []
    tables: list[list[list[str]]] = []
    n_charts = 0
    shape_count = 0
    title_text = title or ""

    for sh in _walk(slide.shapes):
        shape_count += 1
        # tables
        try:
            if sh.has_table:
                tbl = sh.table
                rows: list[list[str]] = []
                for r in tbl.rows:
                    rows.append([c.text.strip() for c in r.cells])
                tables.append(rows)
                continue
        except Exception:
            pass
        # charts
        try:
            if sh.has_chart:
                n_charts += 1
                text_runs.extend(_chart_text(sh))
                continue
        except Exception:
            pass
        # text frames (skip the title; it is captured separately)
        try:
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if not line:
                        # paragraph may hold text without explicit runs
                        line = (para.text or "").strip()
                    if line and line != title_text:
                        text_runs.append(line)
        except Exception:
            continue

    notes = ""
    try:
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
    except Exception:
        notes = ""

    return SlideGT(
        slide_number=number,
        title=title,
        text_runs=text_runs,
        tables=tables,
        notes=notes,
        image_rels=_slide_images(slide),
        n_charts=n_charts,
        shape_count=shape_count,
    )


def extract_ground_truth(pptx_path: str | Path) -> GroundTruth:
    from pptx import Presentation

    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))

    slides = [extract_slide(s, i) for i, s in enumerate(prs.slides, start=1)]

    distinct_images: set[str] = set()
    for s in slides:
        distinct_images.update(s.image_rels)

    gt = GroundTruth(
        deck_slug=deck_slug(pptx_path),
        source_path=str(pptx_path),
        n_slides=len(slides),
        n_notes=sum(1 for s in slides if s.notes.strip()),
        n_tables=sum(len(s.tables) for s in slides),
        n_images=len(distinct_images),
        n_charts=sum(s.n_charts for s in slides),
        slides=slides,
    )
    gt.gt_tokens = count_tokens(gt.all_text() + "\n" + gt.notes_text())
    return gt
