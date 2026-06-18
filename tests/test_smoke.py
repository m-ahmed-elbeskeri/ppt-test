"""Smoke tests — build a tiny synthetic deck (no 12 MB fixture, no COM/network)
and assert the offline path works: availability -> ground truth -> convert ->
metrics -> report, plus the review helpers.

Run:  .venv/Scripts/python -m pytest tests/ -q
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from ppteval import adapters, runner
from ppteval.groundtruth import extract_ground_truth
from ppteval.metrics import compute_metrics
from ppteval.report import build_report


def _make_deck(path) -> None:
    prs = Presentation()
    # Slide 1: title + bullets + speaker notes
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "Quarterly Revenue"
    body = s1.placeholders[1].text_frame
    body.text = "First bullet about EMEA"
    body.add_paragraph().text = "Second bullet about growth"
    s1.notes_slide.notes_text_frame.text = "Speaker note: mention the pipeline"
    # Slide 2: title + a 2x2 table
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "Data Table"
    tbl = s2.shapes.add_table(2, 2, Inches(1), Inches(1.5), Inches(4), Inches(1)).table
    tbl.cell(0, 0).text = "Region"
    tbl.cell(0, 1).text = "Value"
    tbl.cell(1, 0).text = "EMEA"
    tbl.cell(1, 1).text = "123"
    prs.save(str(path))


@pytest.fixture
def deck(tmp_path):
    p = tmp_path / "mini.pptx"
    _make_deck(p)
    return p


def test_availability_lists_all_adapters_without_raising():
    rows = adapters.availability()
    names = {r["name"] for r in rows}
    assert {
        "markitdown", "pptx_custom", "pptx2md", "pymupdf4llm",
        "tika", "docling", "unstructured", "marker",
    } <= names
    avail = {r["name"] for r in rows if r["available"]}
    # offline light adapters must be usable in CI
    assert {"pptx_custom", "markitdown", "pptx2md"} <= avail


def test_ground_truth_extraction(deck):
    gt = extract_ground_truth(deck)
    assert gt.n_slides == 2
    assert gt.n_tables == 1
    assert gt.n_notes == 1
    assert "Quarterly Revenue" in gt.all_text()
    assert "EMEA" in gt.all_text()
    assert "pipeline" in gt.notes_text()


@pytest.mark.parametrize("name", ["pptx_custom", "markitdown", "pptx2md"])
def test_offline_convert_and_metrics(deck, tmp_path, name):
    gt = extract_ground_truth(deck)
    out = adapters.get_adapter(name).convert(deck, tmp_path / name)
    assert out.status == "ok"
    assert out.markdown.strip()
    m = compute_metrics(out, gt)
    assert m["status_ok"] == 1.0
    assert m["out_tokens"] > 0
    assert m["text_recall"] is not None and m["text_recall"] > 0.5


def test_report_builds(deck, tmp_path):
    od = tmp_path / "out"
    od.mkdir()
    gt = runner.run_groundtruth(deck, od)
    outs = runner.run_convert(deck, [adapters.get_adapter("pptx_custom")], od)
    runner.run_metrics(gt, outs, od)
    path = build_report(od)
    assert path.exists()
    txt = path.read_text(encoding="utf-8")
    assert "leaderboard" in txt.lower()
    assert "pptx_custom" in txt
    assert (od / "scorecard.csv").exists()
    assert (od / "results.json").exists()


def test_review_helpers(deck):
    from ppteval.review import _representative_slides, slice_slide

    gt = extract_ground_truth(deck)
    assert len(_representative_slides(gt, 1)) >= 1
    md = "## Slide 1 — A\nfoo\n## Slide 2 — B\nbar"
    assert "foo" in slice_slide(md, 1)
    assert "bar" in slice_slide(md, 2)
    assert slice_slide("no markers here", 1) is None
